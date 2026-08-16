#!/usr/bin/env python3
"""
Sum2Slides Pro — Slide Builder
Generate professional PPTX from structured outline.
Supports multiple templates, brand customization, and auto-layout.
"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(json.dumps({"error": "Missing python-pptx. Install: pip install python-pptx"}))
    sys.exit(1)

import re


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_rr(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_text(slide, left, top, width, height, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    if font_name:
        p.font.name = font_name
    return tf


def build_title_slide(prs, slide_data, template):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    colors = template["style"]["colors"]
    is_dark = template.get("layouts", {}).get("title", {}).get("dark_bg", False)

    if is_dark:
        add_bg(s, hex_to_rgb(colors["primary"]))
        add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
                 hex_to_rgb(colors["accent"]))
        add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
                 slide_data.get("title", "Sum2Slides Summary"),
                 40, hex_to_rgb(colors["text_light"]), True)
        add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.5),
                 slide_data.get("subtitle", ""),
                 20, hex_to_rgb(colors["text_light"]))
        add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.4),
                 slide_data.get("meta", ""),
                 14, RGBColor(0xaa, 0xaa, 0xaa))
    else:
        add_bg(s, hex_to_rgb(colors["background"]))
        add_rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5),
                 hex_to_rgb(colors["primary"]))
        add_text(s, Inches(1.2), Inches(2.0), Inches(10), Inches(1.0),
                 slide_data.get("title", "Sum2Slides Summary"),
                 40, hex_to_rgb(colors["text"]), True)
        add_text(s, Inches(1.2), Inches(3.3), Inches(10), Inches(0.5),
                 slide_data.get("subtitle", ""),
                 20, hex_to_rgb(colors["text_secondary"]))
        add_text(s, Inches(1.2), Inches(4.2), Inches(10), Inches(0.4),
                 slide_data.get("meta", ""),
                 14, hex_to_rgb(colors["text_secondary"]))

    if slide_data.get("brand_logo"):
        try:
            s.shapes.add_picture(slide_data["brand_logo"], Inches(10.5), Inches(0.3),
                                 Inches(2), Inches(1.5))
        except Exception:
            pass


def build_overview_slide(prs, slide_data, template):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    colors = template["style"]["colors"]
    add_bg(s, hex_to_rgb(colors.get("light_bg", "#ffffff")))
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
             hex_to_rgb(colors["primary"]))
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
             "Overview", 32, hex_to_rgb(colors["text"]), True)

    stats = slide_data.get("stats", {})
    stat_items = [
        ("Topics", str(stats.get("topics", 0))),
        ("Decisions", str(stats.get("decisions", 0))),
        ("Action Items", str(stats.get("actions", 0))),
        ("Consensus", str(stats.get("consensus", 0))),
        ("Divergence", str(stats.get("divergence", 0))),
        ("Participants", str(stats.get("speakers", 0))),
        ("Messages", str(stats.get("messages", 0))),
    ]
    for i, (label, value) in enumerate(stat_items):
        col = i % 4
        row = i // 4
        x = Inches(0.8 + col * 3.0)
        y = Inches(1.3 + row * 1.5)
        add_rr(s, x, y, Inches(2.5), Inches(1.2),
               hex_to_rgb(colors["background"]))
        add_text(s, x, y + Inches(0.1), Inches(2.5), Inches(0.6),
                 value, 36, hex_to_rgb(colors["accent"]), True, PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.7), Inches(2.5), Inches(0.4),
                 label, 13, hex_to_rgb(colors["text_secondary"]), align=PP_ALIGN.CENTER)

    if slide_data.get("speakers"):
        add_text(s, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
                 f"Participants: {', '.join(slide_data['speakers'])}",
                 14, hex_to_rgb(colors["text_secondary"]))

    if slide_data.get("duration"):
        add_text(s, Inches(0.8), Inches(6.0), Inches(10), Inches(0.4),
                 f"Session: {slide_data['duration']}",
                 14, hex_to_rgb(colors["text_secondary"]))


def build_topic_slides(prs, topics, template):
    colors = template["style"]["colors"]
    for topic in topics:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s, hex_to_rgb(colors["background"]))
        add_rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5),
                 hex_to_rgb(colors["primary"]))
        add_text(s, Inches(1.0), Inches(0.4), Inches(10), Inches(0.5),
                 topic.get("title", "Topic"), 26,
                 hex_to_rgb(colors["text"]), True)

        # Key points
        y = Inches(1.3)
        points = topic.get("key_points", [])
        for pt in points:
            txt = pt if len(pt) < 120 else pt[:117] + "..."
            add_rr(s, Inches(1.2), y, Inches(10.5), Inches(0.45),
                   hex_to_rgb(colors.get("light_bg", "#f8fafc")))
            add_text(s, Inches(1.5), y + Inches(0.05), Inches(10), Inches(0.35),
                     txt, 13, hex_to_rgb(colors["text"]))
            y += Inches(0.55)
            if y > Inches(6.0):
                break

        # Decisions for topic
        decs = topic.get("decisions", [])
        if decs and y < Inches(5.0):
            add_text(s, Inches(1.0), y + Inches(0.1), Inches(10), Inches(0.3),
                     "Decisions:", 14, hex_to_rgb(colors["accent"]), True)
            y += Inches(0.4)
            for d in decs[:2]:
                txt = d.get("decision", "")[:100]
                add_text(s, Inches(1.3), y, Inches(10), Inches(0.25),
                         f"  {txt}", 12, hex_to_rgb(colors["text_secondary"]))
                y += Inches(0.3)

        # Divergence for topic
        divs = topic.get("divergence", [])
        if divs and y < Inches(5.5):
            add_text(s, Inches(1.0), y + Inches(0.1), Inches(10), Inches(0.3),
                     "Open Questions:", 14, hex_to_rgb(colors.get("warning", colors["accent"])), True)
            y += Inches(0.4)
            for d in divs[:2]:
                txt = d.get("viewpoint", "")[:100]
                add_text(s, Inches(1.3), y, Inches(10), Inches(0.25),
                         f"  {txt}", 12, hex_to_rgb(colors["text_secondary"]))
                y += Inches(0.3)


def build_decision_slide(prs, decisions, template):
    if not decisions:
        return
    colors = template["style"]["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, hex_to_rgb(colors.get("light_bg", "#f8fafc")))
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
             hex_to_rgb(colors["accent"]))
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
             "Key Decisions", 32, hex_to_rgb(colors["text"]), True)

    y = Inches(1.3)
    for dec in decisions:
        add_rr(s, Inches(0.8), y, Inches(11.5), Inches(0.7),
               hex_to_rgb(colors["background"]))
        add_text(s, Inches(1.2), y + Inches(0.05), Inches(10.5), Inches(0.25),
                 dec.get("decision", "")[:150], 14, hex_to_rgb(colors["text"]))
        add_text(s, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.2),
                 f"by {dec.get('by', '')}",
                 11, hex_to_rgb(colors["text_secondary"]))
        y += Inches(0.8)
        if y > Inches(6.5):
            break


def build_consensus_slide(prs, consensus_list, template):
    if not consensus_list:
        return
    colors = template["style"]["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, hex_to_rgb(colors["background"]))
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
             hex_to_rgb(colors.get("success", "#06d6a0")))
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
             "Consensus ✅", 32, hex_to_rgb(colors["text"]), True)

    y = Inches(1.3)
    for c in consensus_list[:6]:
        add_rr(s, Inches(0.8), y, Inches(11.5), Inches(0.7),
               hex_to_rgb(colors.get("light_bg", "#f0fdf4")))
        add_text(s, Inches(1.2), y + Inches(0.05), Inches(10.5), Inches(0.25),
                 c.get("agreement", "")[:150], 14, hex_to_rgb(colors["text"]))
        agreed = c.get("agreed_by", "")
        resp = c.get("respondent_speaker", "")
        if agreed:
            by_text = f"Agreed by {agreed}"
            if resp:
                by_text += f" · in response to {resp}"
            add_text(s, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.2),
                     by_text, 11, hex_to_rgb(colors["text_secondary"]))
        y += Inches(0.8)
        if y > Inches(6.5):
            break


def build_divergence_slide(prs, divergence_list, template):
    if not divergence_list:
        return
    colors = template["style"]["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, hex_to_rgb(colors.get("light_bg", "#fef2f2")))
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
             hex_to_rgb(colors.get("danger", "#ef4444")))
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
             "Open Points 🔄", 32, hex_to_rgb(colors["text"]), True)

    y = Inches(1.5)
    add_text(s, Inches(0.8), y, Inches(11), Inches(0.4),
             "These points need further discussion or alignment:",
             14, hex_to_rgb(colors["text_secondary"]))
    y += Inches(0.5)

    for d in divergence_list[:6]:
        add_rr(s, Inches(0.8), y, Inches(11.5), Inches(0.7),
               hex_to_rgb(colors["background"]))
        add_text(s, Inches(1.2), y + Inches(0.05), Inches(10.5), Inches(0.25),
                 d.get("viewpoint", "")[:150], 14, hex_to_rgb(colors["text"]))
        add_text(s, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.2),
                 f"Raised by {d.get('raised_by', '')}",
                 11, hex_to_rgb(colors["text_secondary"]))
        y += Inches(0.8)
        if y > Inches(6.5):
            break


def build_action_slide(prs, actions, template):
    if not actions:
        return
    colors = template["style"]["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, hex_to_rgb(colors["background"]))
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
             hex_to_rgb(colors["primary"]))
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
             "Action Items", 32, hex_to_rgb(colors["text"]), True)

    y = Inches(1.3)
    add_rr(s, Inches(0.8), y, Inches(11.5), Inches(0.45),
           hex_to_rgb(colors["primary"]))
    add_text(s, Inches(1.0), y + Inches(0.05), Inches(7.0), Inches(0.3),
             "Task", 13, hex_to_rgb(colors["text_light"]), True)
    add_text(s, Inches(8.0), y + Inches(0.05), Inches(2.5), Inches(0.3),
             "Owner", 13, hex_to_rgb(colors["text_light"]), True)

    y += Inches(0.5)
    for action in actions[:8]:
        add_rr(s, Inches(0.8), y, Inches(11.5), Inches(0.45),
               hex_to_rgb(colors.get("light_bg", "#f8fafc")))
        add_text(s, Inches(1.0), y + Inches(0.05), Inches(7.0), Inches(0.3),
                 action.get("task", "")[:100], 12, hex_to_rgb(colors["text"]))
        add_text(s, Inches(8.0), y + Inches(0.05), Inches(2.5), Inches(0.3),
                 action.get("assignee", ""), 12, hex_to_rgb(colors["accent"]), True)
        y += Inches(0.5)
        if y > Inches(6.5):
            break


def build_appendix_slide(prs, extra_data, template):
    colors = template["style"]["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, hex_to_rgb(colors.get("light_bg", "#f8fafc")))
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
             "Appendix", 28, hex_to_rgb(colors["text"]), True)
    add_text(s, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
             extra_data.get("note", "Generated by Sum2Slides Pro"),
             14, hex_to_rgb(colors["text_secondary"]))


def build(outline_data, template_name="default", output_path="output/summary.pptx",
          brand_logo=None, brand_color=None, max_slides=50):
    tmpl_dir = Path(__file__).parent.parent / "templates"
    tmpl_file = tmpl_dir / f"{template_name}.json"
    if not tmpl_file.exists():
        tmpl_file = tmpl_dir / "default.json"

    with open(tmpl_file) as f:
        template = json.load(f)

    if brand_color:
        template["style"]["colors"]["primary"] = brand_color

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    meta = outline_data.get("metadata", {})
    topics = outline_data.get("topics", [])
    decisions = outline_data.get("decisions", [])
    actions = outline_data.get("action_items", [])
    consensus = outline_data.get("consensus", [])
    divergence = outline_data.get("divergence", [])

    # Title slide
    subtitle_parts = []
    tc = len(topics)
    if tc: subtitle_parts.append(f"{tc} topic{'s' if tc != 1 else ''}")
    dc = meta.get("decisions_found", len(decisions))
    if dc: subtitle_parts.append(f"{dc} decision{'s' if dc != 1 else ''}")
    ac = meta.get("action_items_found", len(actions))
    if ac: subtitle_parts.append(f"{ac} action{'s' if ac != 1 else ''}")
    conc = meta.get("consensus_found", len(consensus))
    if conc: subtitle_parts.append(f"{conc} consensus{' point' if conc != 1 else ''}")
    divc = meta.get("divergence_found", len(divergence))
    if divc: subtitle_parts.append(f"{divc} open point{'s' if divc != 1 else ''}")

    build_title_slide(prs, {
        "title": "Sum2Slides Summary",
        "subtitle": " · ".join(subtitle_parts) if subtitle_parts else "Brainstorm Summary",
        "meta": f"Participants: {', '.join(meta.get('speakers', []))}",
        "brand_logo": brand_logo
    }, template)

    # Overview slide
    build_overview_slide(prs, {
        "stats": {
            "topics": meta.get("topics_found", len(topics)),
            "decisions": meta.get("decisions_found", len(decisions)),
            "actions": meta.get("action_items_found", len(actions)),
            "consensus": meta.get("consensus_found", len(consensus)),
            "divergence": meta.get("divergence_found", len(divergence)),
            "speakers": len(meta.get("speakers", [])),
            "messages": meta.get("total_messages", 0)
        },
        "speakers": meta.get("speakers", [])
    }, template)

    # Topic slides
    slide_count = 2
    for topic in topics:
        if slide_count >= max_slides:
            break
        build_topic_slides(prs, [topic], template)
        slide_count += 1

    # Consensus slide
    if consensus and slide_count < max_slides:
        build_consensus_slide(prs, consensus[:6], template)
        slide_count += 1

    # Divergence slide
    if divergence and slide_count < max_slides:
        build_divergence_slide(prs, divergence[:6], template)
        slide_count += 1

    # Decision slide
    if decisions and slide_count < max_slides:
        build_decision_slide(prs, decisions[:6], template)
        slide_count += 1

    # Action slide
    if actions and slide_count < max_slides:
        build_action_slide(prs, actions[:8], template)
        slide_count += 1

    # Appendix
    if slide_count < max_slides:
        build_appendix_slide(prs, {
            "note": "Generated by Sum2Slides Pro · cqdev-ai / ChengQian"
        }, template)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path, slide_count


def main():
    parser = argparse.ArgumentParser(description="Sum2Slides Pro — Slide Builder")
    parser.add_argument("--outline", required=True, help="Outline JSON from extractor")
    parser.add_argument("--template", default="meeting",
                        choices=["default", "meeting", "brainstorm", "retro"],
                        help="Slide template")
    parser.add_argument("--output", default="output/summary.pptx", help="Output PPTX path")
    parser.add_argument("--brand-logo", help="Company logo image")
    parser.add_argument("--brand-color", help="Theme accent hex color")
    parser.add_argument("--max-slides", type=int, default=50, help="Max slides")
    args = parser.parse_args()

    with open(args.outline) as f:
        outline_data = json.load(f)

    output, count = build(outline_data, args.template, args.output,
                          args.brand_logo, args.brand_color, args.max_slides)

    print(json.dumps({
        "output": output,
        "slides": count,
        "template": args.template,
        "brand_color": args.brand_color or "default"
    }, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
